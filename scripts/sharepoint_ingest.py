#!/usr/bin/env python3
"""
SharePoint RAG Ingest — Delegated Permissions (Device Code Flow)

Connects to SharePoint via Microsoft Graph API using delegated auth.
First run: prints a URL + code for browser login (device code flow).
Subsequent runs: uses cached refresh token automatically (~90 days).

Features:
- Client name masking (specific + pattern-based)
- Person name masking
- Phase-aware chunking (Prepare / Explore / Realize / Deploy / Run)
- Recursive subfolder traversal
- Structured chunk metadata (phase, client, source)

Usage:
    python3.11 scripts/sharepoint_ingest.py           # Graph API mode
    python3.11 scripts/sharepoint_ingest.py --local   # local raw/ folder mode (POC)

Env vars:
    GRAPH_TENANT_ID       Azure AD Directory (tenant) ID
    GRAPH_CLIENT_ID       App registration Application (client) ID
    GRAPH_CLIENT_SECRET   Client secret value
    SHAREPOINT_SITE_URL   e.g. https://ts.accenture.com/sites/S4_HANA_POD_Harvesting
    SHAREPOINT_LIBRARY    Document library name (default: Shared Documents)
    SHAREPOINT_SUBFOLDER  Subfolder path within the library

Install:
    pip3.11 install msal requests python-docx pymupdf python-pptx openpyxl
    # Proper NER masking (recommended — else falls back to regex name masking):
    pip3.11 install spacy && python3.11 -m spacy download en_core_web_lg
"""

import os
import sys
import json
import shutil
import hashlib
import logging
import re
import argparse
from pathlib import Path
from datetime import datetime

# ── CONFIG ───────────────────────────────────────────────────────────────────
TENANT_ID  = os.environ.get("GRAPH_TENANT_ID", "")
CLIENT_ID  = os.environ.get("GRAPH_CLIENT_ID", "")
SITE_URL   = os.environ.get("SHAREPOINT_SITE_URL", "")
LIBRARY    = os.environ.get("SHAREPOINT_LIBRARY", "Shared Documents")
SUBFOLDER  = os.environ.get("SHAREPOINT_SUBFOLDER", "")

SCOPES     = ["Sites.Read.All", "Files.Read.All"]
GRAPH_BASE = "https://graph.microsoft.com/v1.0"

BASE_DIR   = Path(__file__).resolve().parent.parent
BRAIN_DIR  = BASE_DIR / "brain" / "sharepoint"
RAW_DIR    = BRAIN_DIR / "raw"
CHUNKS_DIR = BRAIN_DIR / "chunks"
TOKEN_CACHE = BASE_DIR / "brain" / ".token_cache.json"
LOG_FILE   = BASE_DIR / "brain" / "ingest.log"

CHUNK_WORDS   = 512
CHUNK_OVERLAP = 64
SUPPORTED_EXT = {".docx", ".pdf", ".pptx", ".txt", ".md", ".xlsx"}
_CONTENT_SCAN = 20_000        # chars of doc text scanned for phase/role/deliverable hints

# ── PHASES ────────────────────────────────────────────────────────────────────
# SAP Activate: Discover / Prepare / Explore / Realize / Deploy / Run
PHASES = ["discover", "prepare", "explore", "realize", "deploy", "run"]

# Phase keyword hints — SAP Activate roadmap workstreams + official phase-activity
# definitions (Discover/Prepare/Explore/Realize/Deploy/Run).
_PHASE_HINTS = [
    ("discover", [
        "discover", "cloud trial", "trial system", "cloud mindset",
        "discovery assessment", "digital discovery assessment", "dda",
        "value discovery", "application value", "value scoping", "scoping",
        "solution scope", "target solution model", "user enablement",
    ]),
    ("prepare",  [
        "getting started", "onboarding", "project initiation", "project governance",
        "project standards", "project charter", "project goals", "project plan",
        "kick off", "ko deck", "executive sponsorship", "business value objective",
        "implementation strategy", "upgrade strategy", "roles and responsibilities",
        "team self-enablement", "self-enablement", "tool access", "sap cloud alm",
        "sap cbc", "starter system", "cloud starter system", "initial access",
        "business driven configuration", "configuration assessment",
        "fit-to-standard analysis preparation", "system preparation",
        "analytics approach", "analytics plan", "data approach", "data plan",
        "integration setup", "enablement strategy", "organizational change",
        "raci", "sow", "statement of work", "roles matrix", "authorization matrix",
        "l4 plan", "sterco", "project reporting", "project tracking",
    ]),
    ("explore",  [
        "fit-to-standard", "fit to standard", "fit-to-standard analysis", "f2s",
        "f2s deck", "fit-to-standard documentation", "delta requirement",
        "delta configuration", "configuration definition", "backlog",
        "identity and access management", "integration prerequisites",
        "extension planning", "solution extension preparation",
        "integration planning", "integration design", "test planning",
        "data load preparation", "learning needs analysis", "phase closure",
        "analytics planning", "analytics design", "solution design workshop",
        "show-and-tell", "show and tell", "kdd", "workshop", "business process",
        "wricef", "architecture design", "solution confirmation", "user stor",
        "customer execution", "standard processes",
    ]),
    ("realize",  [
        "required configuration", "solution configuration", "incremental build",
        "time-boxed iteration", "time boxed iteration", "iteration",
        "development system", "test environment", "replicate the solution",
        "integrated business", "end-to-end testing", "end to end testing",
        "key user training", "production environment", "end user training",
        "activate new scope", "solution extension", "extension development",
        "extension deployment", "test preparation", "test execution",
        "sprint planning", "sprint execution", "data migration development",
        "data migration test", "support operations", "handover plan",
        "enablement content", "analytics configuration", "functional design",
        " fd ", "fd -", "fd sample", "technical design", "tdd", " td ",
        "rap", "cap code", "ui code", "form wizard", "interface", "iflow",
        "data strategy", "data cleansing", "test strategy", "test case",
        "test script", "test data", "defect resolver", "code quality",
        "technical strategy",
    ]),
    ("deploy",   [
        "system go-live", "go-live", "go live", "production go-live",
        "cutover preparation", "production cutover", "cutover activities",
        "cutover plan", "transition plan", "organization readiness",
        "organizational change management", "ocm", "post-go-live support",
        "post go live support", "switch business operations", "system tests",
        "release update cycles", "operations readiness", "change impact",
        "change strategy", "training material", "kut", "eut",
        "communication template", "cutover", "copy reference", "talent agent",
        "deployment",
    ]),
    ("run",      [
        "ongoing operations", "ongoing system operations", "operability",
        "system availability", "performance levels", "continuous improvement",
        "continuous change management", "continuous learning", "release cycles",
        "value management", "release update", "activate new scope", "run support",
        "incident", "service request", "release impact", "knowledge steward",
        "autonomous ops",
    ]),
]

def detect_phase(path_str: str, text: str = "") -> str:
    p = path_str.lower().replace("\\", "/")
    # Folder name takes priority — matches plain (Realize/), numbered (4.Realize/),
    # or prefixed (5.Deploy/) folder names anywhere in the path.
    for phase in PHASES:
        if re.search(rf'(?:^|/)\d*\.?{phase}(?:/|$)', p):
            return phase.capitalize()
    # Fallback: keyword hints over the filename AND the document content, so a
    # generically-named doc (e.g. "Glossary.xlsx") is classified by what's inside.
    hay = p + "\n" + text[:_CONTENT_SCAN].lower()
    for phase, hints in _PHASE_HINTS:
        if any(h in hay for h in hints):
            return phase.capitalize()
    return "General"

# ── AGENT ROLES ───────────────────────────────────────────────────────────────
# Each entry: (agent_role_key, [keywords to match in path or filename])
# ── SAP STANDARD BPD DOCS ─────────────────────────────────────────────────────
# SAP-delivered Business Process Documentation, named <SCOPE>_S4CLD<ver>_BPD_...
# (e.g. 1MR_S4CLD2402_BPD_EN_US.docx). These are SAP STANDARD reference content —
# separated from client delivery docs (source_system=sap_bpd, phase=Reference) and
# linked back to the scope catalog by the scope item ID in the filename.
_BPD_RE = re.compile(r"^([0-9A-Z]{2,4})_S4CLD\d+_BPD", re.IGNORECASE)

def detect_sap_bpd(filename: str):
    """Return the scope item ID if the file is an SAP-standard BPD, else None."""
    m = _BPD_RE.match(filename.strip())
    return m.group(1).upper() if m else None

_AGENT_ROLE_KEYWORDS = [
    ("pmo_agent", [
        "pmo_agent", "pmo agent", "project charter", "kick off", "ko deck",
        "l4 plan", "l4plan", "onboarding kit", "onboarding", "sterco",
        "raci", "sow", "statement of work", "project sow",
    ]),
    ("security_agent", [
        "security_agent", "security agent", "orion agent", "roles matrix",
        "authorization matrix", "role requirement", "roles & author",
        "identity and access", "iam", "cloud trial",
    ]),
    ("solution_confirmation_agent", [
        "solution_confirmation_agent", "solution confirmation",
        "business process", "fit to standard", "f2s", "f2s deck",
        "kdd", "digital discovery", "dda", "workshop", "wricef",
        "architecture design", "user stor", "bdcq",
    ]),
    ("functional_agent", [
        "functional_agent", "functional agent",
        "functional design", " fd ", "fd -", "fd sample", "fd_",
        "config wizard", "config rationale",
    ]),
    ("build_agent", [
        "build_agent", "build agent",
        "technical design", "tdd", " td ", "td ref", "td_",
        "rap code", "cap code", "ui code", "form wizard",
        "interface functional", "interface spec", "interface mapper",
        "iflow", "integration iflow", "technical strategy",
        "code quality", "impact analysis",
    ]),
    ("data_agent", [
        "data_agent", "data agent",
        "data migration", "data strategy", "data cleansing",
        "data enrichment", "data profiler", "adcmc", "migration approach",
    ]),
    ("qe_agent", [
        "qe_agent", "qe agent", "quality engineering",
        "test strategy", "test case", "test script", "test data",
        "test plan", "defect resolver", "defect resolution",
    ]),
    ("change_talent_agent", [
        "change_talent_agent", "change talent",
        "change impact", "change strategy", "training material",
        "training plan", "kut ", "eut ", "communication template",
        "talent agent",
    ]),
    ("deployment_agent", [
        "deployment_agent", "deployment agent",
        "cutover", "cutover plan", "copy reference", "task generator",
    ]),
    ("run_support_agent", [
        "run_support_agent", "run support",
        "defect resolver", "incident resolution", "service request",
        "release impact", "knowledge steward", "autonomous ops",
    ]),
]

def detect_agent_role(path_str: str, text: str = "") -> str:
    hay = path_str.lower() + "\n" + text[:_CONTENT_SCAN].lower()
    for role_key, keywords in _AGENT_ROLE_KEYWORDS:
        if any(kw in hay for kw in keywords):
            return role_key
    return "general"

# ── DELIVERABLE TYPES ─────────────────────────────────────────────────────────
# Each entry: (deliverable_type, [keywords to match])
_DELIVERABLE_KEYWORDS = [
    ("project_charter",          ["project charter"]),
    ("kickoff_deck",             ["kick off", "ko deck", "kickoff"]),
    ("project_plan",             ["l4 plan", "l4plan", "project plan", "milestone"]),
    ("onboarding_kit",           ["onboarding kit", "onboarding"]),
    ("raci_matrix",              ["raci"]),
    ("statement_of_work",        ["sow", "statement of work", "project sow"]),
    ("roles_authorization",      ["roles matrix", "role requirement", "authorization matrix",
                                  "roles & author"]),
    ("discovery_assessment",     ["dda", "digital discovery", "discovery assessment"]),
    ("business_process_design",  ["business process design", "business process"]),
    ("fit_to_standard",          ["fit to standard", "f2s deck", "f2s ", "f2sdeck"]),
    ("kdd",                      ["kdd deck", "kdd "]),
    ("workshop_analysis",        ["workshop analysis", "design workshop", "workshop deck"]),
    ("wricef_inventory",         ["wricef"]),
    ("architecture_design",      ["architecture design"]),
    ("functional_design",        ["functional design", " fd ", "fd -", "fd sample", "fd_"]),
    ("configuration",            ["config wizard", "config rationale", "configuration"]),
    ("technical_design",         ["technical design", "tdd ref", " tdd ", " td ref"]),
    ("rap_code",                 ["rap code", "rap template"]),
    ("cap_code",                 ["cap code", "cap template"]),
    ("ui_code",                  ["ui code", "ui template", "fiori"]),
    ("form_wizard",              ["form wizard", "adobe form", "form template"]),
    ("interface_spec",           ["interface functional", "interface spec", "interface mapper",
                                  "interface template"]),
    ("integration_iflow",        ["iflow", "integration iflow", "integration flow"]),
    ("technical_strategy",       ["technical strategy"]),
    ("data_strategy",            ["data strategy"]),
    ("data_migration",           ["data migration", "migration approach", "migration agent"]),
    ("data_profiler",            ["data profiler", "adcmc"]),
    ("data_cleansing",           ["data cleansing", "data enrichment"]),
    ("test_strategy",            ["test strategy"]),
    ("test_cases",               ["test case", "test script", "test scenario"]),
    ("test_data",                ["test data"]),
    ("change_impact",            ["change impact"]),
    ("change_strategy",          ["change strategy"]),
    ("training_material",        ["training material", "training plan", "kut ", "eut "]),
    ("communication_template",   ["communication template"]),
    ("cutover_plan",             ["cutover plan", "cutover approach"]),
    ("copy_reference",           ["copy reference"]),
    ("defect_resolution",        ["defect resolver", "defect resolution"]),
    ("knowledge_base",           ["knowledge steward", "knowledge base"]),
    ("incident_resolution",      ["incident resolution", "service request"]),
    ("release_impact",           ["release impact"]),
]

def detect_deliverable_type(path_str: str, text: str = "") -> str:
    hay = path_str.lower() + "\n" + text[:_CONTENT_SCAN].lower()
    for deliverable, keywords in _DELIVERABLE_KEYWORDS:
        if any(kw in hay for kw in keywords):
            return deliverable
    return "reference_document"

# ── CONTENT TYPE ──────────────────────────────────────────────────────────────
_CONTENT_TYPE_KEYWORDS = {
    "template":    ["template", "templ", "blank", "format"],
    "example":     ["sample", "example", "demo", "ver 1", "ver1", "ver 2", "ver2"],
    "reference":   ["reference", "ref ", "guide", "handbook", "playbook"],
    "methodology": ["approach", "strategy", "framework", "methodology", "process"],
    "assessment":  ["assessment", "results", "discovery", "analysis", "dda"],
}

def detect_content_type(filename: str) -> str:
    f = filename.lower()
    for ctype, keywords in _CONTENT_TYPE_KEYWORDS.items():
        if any(kw in f for kw in keywords):
            return ctype
    return "document"

# ── CLIENT NAMES ─────────────────────────────────────────────────────────────
KNOWN_CLIENTS = [
    "BOBST", "CAMPARI", "CUMMINS", "BUMA", "MARS",
    "Altor Damas", "CDI", "AXA",
]

def detect_client_from_path(path_str: str) -> str:
    for client in KNOWN_CLIENTS:
        if client.lower() in path_str.lower():
            return client
    return "Unknown"

# ── LOGGING ───────────────────────────────────────────────────────────────────
LOG_FILE.parent.mkdir(parents=True, exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(), logging.FileHandler(LOG_FILE)],
)
log = logging.getLogger("sharepoint_ingest")

# ── MASKING ───────────────────────────────────────────────────────────────────
# Build client name pattern from known list
_client_alternatives = "|".join(re.escape(c) for c in KNOWN_CLIENTS)

# SAP / business vocabulary — capitalised words that are NOT person names.
# Two jobs: (1) guard the regex standalone-name fallback, (2) filter spaCy NER
# so it never masks a business term like "Financial Close" as a PERSON/ORG.
_BUSINESS_VOCAB = {
    # finance / controlling
    "financial", "finance", "close", "closing", "accounting", "account",
    "general", "ledger", "asset", "controlling", "treasury", "banking",
    "payment", "receivable", "payable", "fixed", "depreciation", "revenue",
    "margin", "profit", "cost", "center", "centre", "profitability", "budget",
    "actuals", "settlement", "allocation", "reconciliation", "dunning",
    "credit", "debit", "journal", "posting", "clearing", "consolidation",
    "group", "gaap", "ifrs", "currency", "exchange", "valuation", "accrual",
    # sales / procurement / supply
    "sales", "order", "purchase", "purchasing", "procurement", "sourcing",
    "supplier", "vendor", "customer", "quotation", "contract", "requisition",
    "delivery", "shipment", "shipping", "goods", "receipt", "issue", "billing",
    "invoice", "pricing", "condition", "rebate", "returns", "fulfillment",
    "supply", "chain", "logistics", "inventory", "warehouse", "stock",
    "replenishment", "transfer", "transportation", "handling", "picking",
    "packing", "putaway", "batch", "serial", "availability", "promising",
    # manufacturing / plm / am
    "manufacturing", "production", "planning", "process", "discrete", "repetitive",
    "shop", "floor", "routing", "bom", "work", "resource", "capacity", "mrp",
    "maintenance", "asset", "corrective", "preventive", "engineering",
    "product", "development", "variant", "configuration", "quality",
    "inspection", "notification", "defect", "calibration",
    # org / master data / cross
    "master", "data", "material", "business", "partner", "company", "code",
    "plant", "organization", "organizational", "unit", "hierarchy", "cross",
    "management", "system", "responsibility", "situation", "operations",
    "operational", "compliance", "governance", "risk", "audit", "tax", "legal",
    "third", "party", "enterprise", "sector", "retail", "baseline", "accelerator",
    "country", "region", "global", "local", "central", "event", "condition",
    # sap platform / clean core / extensibility
    "sap", "s4hana", "s4", "hana", "fiori", "launchpad", "activate", "clean",
    "core", "extensibility", "adaptation", "adapt", "workflow", "flexible",
    "output", "determination", "custom", "field", "logic", "released", "object",
    "namespace", "transport", "tenant", "badi", "cds", "rap", "abap", "odata",
    "api", "btp", "cap", "steampunk", "keyuser", "developer", "sidebyside",
    "analytical", "query", "app", "application", "standard", "scope", "item",
    "solution", "cloud", "public", "edition", "private", "onpremise",
    # activate phases / delivery
    "discover", "prepare", "explore", "realize", "deploy", "run", "phase",
    "sprint", "cutover", "onboarding", "enablement", "adoption", "value",
    "test", "testing", "unit", "integration", "acceptance", "regression",
    "scenario", "script", "case", "change", "impact", "training", "communication",
    "deployment", "release", "update", "cycle", "support", "incident", "handover",
    "readiness", "hypercare", "golive", "mock", "dress", "rehearsal",
    # documents / delivery artifacts
    "functional", "technical", "design", "specification", "charter", "matrix",
    "inventory", "roadmap", "assessment", "discovery", "analysis", "approach",
    "plan", "guide", "guideline", "kickoff", "kick", "off", "deck", "status",
    "power", "workshop", "requirement", "requirements", "template", "reference",
    "document", "documentation", "report", "interface", "migration", "strategy",
    "enrichment", "cleansing", "resolver", "wricef", "ricefw", "raci", "sow",
    "fit", "gap", "blueprint", "backlog", "story", "epic", "feature",
    # geography words spaCy may tag as GPE/ORG in headings
    "north", "south", "east", "west", "america", "americas", "europe", "asia",
    "pacific", "emea", "apac", "united", "states", "kingdom",
}

# Vendors / products / partners that spaCy will tag as ORG but must NEVER be
# masked — masking these would destroy the SAP knowledge in the brain.
_NEVER_MASK_ORGS = {
    "sap", "s/4hana", "s4hana", "s/4", "fiori", "hana", "ariba", "fieldglass",
    "concur", "successfactors", "sap successfactors", "sap ariba", "sap fieldglass",
    "opentext", "vistex", "microsoft", "microsoft graph", "azure", "amazon",
    "aws", "amazon web services", "bedrock", "google", "oracle", "accenture",
    "ibm", "salesforce", "workday", "servicenow", "adobe", "openai", "anthropic",
    "claude", "titan", "github", "gitlab", "jira", "confluence", "sharepoint",
    "teams", "outlook", "excel", "powerpoint", "word", "python", "node", "nodejs",
    "javascript", "ui5", "sapui5", "cap", "btp", "cpi",
}

# Regex fallback (used only when spaCy is unavailable): standalone Firstname
# Lastname, guarded by the business-vocab allowlist.
_STANDALONE_NAME_RE = re.compile(
    r"\b[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})?\b"
)

def _mask_standalone_name(m: "re.Match") -> str:
    phrase = m.group(0)
    if any(w.lower() in _BUSINESS_VOCAB for w in phrase.split()):
        return phrase          # looks like an SAP/business term — leave it
    return "[PERSON]"

_MASK_RULES = [
    # ── 1. CREDENTIALS (highest risk — first) ─────────────────────────────────
    (re.compile(
        r"(?i)(?:password|passwd|api[_\s]?key|access[_\s]?key|secret[_\s]?key"
        r"|token|bearer|credential|auth[_\s]?key)\s*[:=]\s*\S+",
    ), "[CREDENTIAL]"),

    # ── 2. STRUCTURED IDENTIFIERS (specific formats — before greedy rules) ─────
    # Email addresses
    (re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"), "[EMAIL]"),
    # Internal/client URLs (not public SAP docs)
    (re.compile(
        r"https?://(?!help\.sap\.com|api\.sap\.com|cap\.cloud\.sap|ui5\.sap\.com"
        r"|www\.sap\.com|discovery\.sap\.com)[^\s\"'<>]+"
    ), "[INTERNAL_URL]"),
    # SAP tenant URLs (myXXXXXX.s4hana.ondemand.com)
    (re.compile(r"\bmy[A-Za-z0-9]+\.s4hana\.ondemand\.com\b"), "[SAP_TENANT_URL]"),
    # IP addresses (v4) — before phone, so phone can't eat the digit groups
    (re.compile(r"\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b"), "[IP_ADDRESS]"),
    # SAP transport requests (e.g. NPLK900123) — before phone
    (re.compile(r"\b[A-Z]{3}[KO]\d{6}\b"), "[TRANSPORT]"),
    # SAP logical system names (client abbrev + CLNT + client no.) — before phone
    (re.compile(r"\b[A-Z]{2,10}CLNT\d{3}\b"), "[LOGICAL_SYSTEM]"),
    # Employee IDs (Accenture I/C format) — before phone
    (re.compile(r"\b[IC]\d{6,7}\b"), "[EMP_ID]"),
    # Project ticket IDs (PROJ-1234)
    (re.compile(r"\b[A-Z]{2,6}-\d{3,}\b"), "[TICKET]"),
    # PO / contract numbers
    (re.compile(r"\b(?:PO|CONTRACT|ORDER)[-\s]?\d{4,}\b", re.IGNORECASE), "[CONTRACT_REF]"),

    # ── 3. CLIENT / COMPANY NAMES ─────────────────────────────────────────────
    # SAP namespace objects named after client (ZBOBST_, ZCDI_, …) — before client
    (re.compile(rf"\bZ(?:{_client_alternatives})[_A-Z0-9]*\b", re.IGNORECASE), "[CLIENT_OBJECT]"),
    # Known client names (exact, case-insensitive)
    (re.compile(rf"\b(?:{_client_alternatives})\b", re.IGNORECASE), "[CLIENT]"),
    # Generic company names (Siemens AG, Bosch GmbH, Contoso Industries, …)
    # Legal-form abbreviations + strong full-word company suffixes that rarely
    # collide with SAP terminology in a "Capword <suffix>" shape.
    (re.compile(
        r"\b[A-Z][A-Za-z&]+(?:\s+[A-Z][A-Za-z&]+)*\s+"
        r"(?:AG|GmbH|Ltd|Inc|Corp|SE|NV|PLC|SA|LLC|LLP|BV|SAS|SpA"
        r"|Corporation|Incorporated|Limited|Industries|Holdings|Enterprises)\b"
    ), "[CLIENT]"),

    # ── 4. FINANCIAL (rate before amount, so "$1,800/day" → RATE) ─────────────
    (re.compile(
        r"(?:USD|EUR|GBP|€|\$|£)\s?\d[\d,\.]+\s?(?:/\s?day|per\s+day)",
        re.IGNORECASE
    ), "[RATE]"),
    (re.compile(
        r"(?:USD|EUR|GBP|CHF|€|\$|£)\s?\d[\d,\.]*\s?(?:K|M|B|thousand|million)?",
        re.IGNORECASE
    ), "[AMOUNT]"),

    # ── 5. PROJECT CODENAMES (before the greedy standalone-name rule) ─────────
    (re.compile(r"\bProject\s+[A-Z][A-Za-z]+\b"), "[PROJECT]"),

    # ── 6. PERSON NAMES ───────────────────────────────────────────────────────
    # Titled names (Mr/Mrs/Dr/Prof …)
    (re.compile(
        r"\b(?:Mr|Mrs|Ms|Miss|Dr|Prof|Eng)\.?\s+[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b"
    ), "[PERSON]"),
    # Author/contact/owner fields
    (re.compile(
        r"(?i)(?:author|by|prepared by|created by|modified by|contact|owner|lead"
        r"|reviewed by|approved by|assigned to)\s*:\s*([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)"
    ), "[PERSON]"),
    # NOTE: standalone Firstname-Lastname is handled by the spaCy NER pass in
    # mask() (or the regex fallback _STANDALONE_NAME_RE if spaCy is unavailable).

    # ── 7. PHONE NUMBERS (greediest digits — LAST; needs + or phone keyword) ──
    # International: leading +country code
    (re.compile(r"\+\d{1,3}[\s\-.]?\(?\d{1,4}\)?(?:[\s\-.]?\d{2,4}){2,5}\b"), "[PHONE]"),
    # Contextual: preceded by a phone/tel/mobile/fax label
    (re.compile(
        r"(?i)(?:phone|tel|telephone|mobile|cell|fax|call)\s*[:.]?\s*"
        r"\+?[\d][\d\s\-.()]{6,}\d"
    ), "[PHONE]"),
]

# ── NER (proper person/org detection via spaCy) ───────────────────────────────
# Loaded lazily and cached. Set SPACY_MODEL to override; en_core_web_lg gives the
# best CPU accuracy without transformers. Only the NER pipe is kept (speed).
_NER = "unset"   # sentinel until first load attempt

def _get_ner():
    global _NER
    if _NER != "unset":
        return _NER
    try:
        import spacy
    except ImportError:
        log.warning("spaCy not installed — using regex name fallback. For proper "
                    "NER: pip3.11 install spacy && python3.11 -m spacy download en_core_web_lg")
        _NER = None
        return _NER
    override = os.environ.get("SPACY_MODEL", "").strip()
    candidates = [override] if override else ["en_core_web_lg", "en_core_web_sm"]
    for name in candidates:
        if not name:
            continue
        try:
            _NER = spacy.load(name, disable=["parser", "lemmatizer", "tagger",
                                             "attribute_ruler"])
            # NER-only pipeline is light on memory — allow large documents.
            _NER.max_length = 3_000_000
            log.info("NER model loaded: %s", name)
            return _NER
        except OSError:
            continue
    log.warning("No spaCy model found — using regex name fallback. Install one: "
                "python3.11 -m spacy download en_core_web_lg")
    _NER = None
    return _NER

# Placeholder labels already inserted by the regex pass — the NER pass must never
# re-mask these (e.g. spaCy tagging "EMP_ID" inside "[EMP_ID]" as an ORG).
_MASK_LABELS = {
    "CREDENTIAL", "EMAIL", "INTERNAL_URL", "SAP_TENANT_URL", "IP_ADDRESS",
    "TRANSPORT", "LOGICAL_SYSTEM", "EMP_ID", "TICKET", "CONTRACT_REF",
    "CLIENT_OBJECT", "CLIENT", "RATE", "AMOUNT", "PROJECT", "PERSON", "PHONE",
}

# NER runs in windows so a single huge document can't exhaust memory (spaCy needs
# ~1GB per 100k chars). Windows break at newlines to avoid splitting an entity.
_NER_WINDOW = 100_000

def _ner_mask(text: str, nlp) -> str:
    """Mask names/orgs via spaCy NER, windowing long text to bound memory."""
    if len(text) <= _NER_WINDOW:
        return _ner_mask_segment(text, nlp)
    out, i, n = [], 0, len(text)
    while i < n:
        end = min(i + _NER_WINDOW, n)
        if end < n:                              # extend to next newline if close
            nl = text.find("\n", end)
            if nl != -1 and nl - end < 5000:
                end = nl + 1
        out.append(_ner_mask_segment(text[i:end], nlp))
        i = end
    return "".join(out)

def _ner_mask_segment(text: str, nlp) -> str:
    """Mask PERSON (and unknown-client ORG) entities in one NER-sized segment."""
    doc = nlp(text)
    spans = []
    for ent in doc.ents:
        etext = ent.text.strip()
        if not etext or etext.startswith("["):      # already a mask token
            continue
        # skip our own placeholder labels (e.g. "EMP_ID" inside "[EMP_ID]")
        if etext.strip("[]").upper() in _MASK_LABELS:
            continue
        toks = re.findall(r"[A-Za-z][A-Za-z&/.\-]*", etext)
        low  = etext.lower()
        if any(t.upper() in _MASK_LABELS for t in toks):
            continue
        # never touch SAP/business vocabulary or known vendors/products
        if any(t.lower() in _BUSINESS_VOCAB for t in toks):
            continue
        if low in _NEVER_MASK_ORGS or any(t.lower() in _NEVER_MASK_ORGS for t in toks):
            continue
        if ent.label_ == "PERSON":
            repl = "[PERSON]"
        elif ent.label_ == "ORG":
            # NER catches unknown client orgs the regex misses — but require a
            # MULTI-WORD name so common single words that happen to be companies
            # (e.g. "Reach", "Order") are never masked. Single-word clients belong
            # in KNOWN_CLIENTS (exact match) instead.
            if " " not in etext or not any(c.isupper() for c in etext):
                continue
            repl = "[CLIENT]"
        else:
            continue                                 # ignore GPE/DATE/etc.
        spans.append((ent.start_char, ent.end_char, repl))
    # apply right-to-left so offsets stay valid
    for start, end, repl in sorted(spans, key=lambda s: s[0], reverse=True):
        text = text[:start] + repl + text[end:]
    return text

def mask(text: str) -> str:
    """Hybrid masking: regex for structured PII, spaCy NER for names/orgs."""
    # 1. structured + known-entity regex rules
    for pattern, replacement in _MASK_RULES:
        text = pattern.sub(replacement, text)
    # 2. proper NER pass for person/org names (fallback to regex if unavailable)
    nlp = _get_ner()
    if nlp is not None:
        text = _ner_mask(text, nlp)
    else:
        text = _STANDALONE_NAME_RE.sub(_mask_standalone_name, text)
    return text

# ── TEXT EXTRACTION ───────────────────────────────────────────────────────────
def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext in (".txt", ".md"):
            return path.read_text(errors="ignore")
        if ext == ".docx":
            import docx
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if ext == ".pdf":
            import fitz
            doc = fitz.open(path)
            return "\n".join(page.get_text() for page in doc)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join(str(c) for c in row if c is not None)
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts)
    except Exception as e:
        log.warning("Extraction failed [%s]: %s", path.name, e)
    return ""

# ── CHUNKING ──────────────────────────────────────────────────────────────────
def chunk(text: str) -> list:
    words, chunks, i = text.split(), [], 0
    while i < len(words):
        c = " ".join(words[i:i + CHUNK_WORDS])
        if c.strip():
            chunks.append(c)
        i += CHUNK_WORDS - CHUNK_OVERLAP
    return chunks

def _safe_str(s: str) -> str:
    """Repair a filename/path carrying non-UTF-8 bytes (surrogate escapes from
    Windows/Mac encodings — e.g. an en-dash saved as byte 0x96) so it can be
    logged, stored in JSON, and read back as clean UTF-8 without a
    UnicodeEncodeError. Undecodable bytes collapse to '?'."""
    return s.encode("utf-8", "replace").decode("utf-8")


def _ingest_one_local(f) -> int:
    """Ingest a single local file into chunk JSONs. Returns the chunk count.
    Raising is fine — the caller skips the file and keeps going."""
    rel_path_raw   = str(f.relative_to(RAW_DIR))   # may hold surrogate bytes
    rel_path       = _safe_str(rel_path_raw)        # clean UTF-8 for logs/JSON
    safe_name      = _safe_str(f.name)

    # Extract + mask first, so classification can read the document content
    # (not just the filename) — many delivery docs have generic names.
    text   = extract_text(f)
    text   = mask(text)

    bpd_scope      = detect_sap_bpd(safe_name)
    if bpd_scope:                       # SAP standard BPD — reference, not delivery
        source_system, phase, agent_role = "sap_bpd", "Reference", "reference"
        deliverable, scope_item_id = "business_process_doc", bpd_scope
    else:                               # client delivery document
        source_system, scope_item_id = "sharepoint", None
        phase       = detect_phase(rel_path, text)
        agent_role  = detect_agent_role(rel_path, text)
        deliverable = detect_deliverable_type(rel_path, text)
    content_type   = detect_content_type(safe_name)

    log.info("Processing: %s [src=%s, phase=%s, agent=%s, deliverable=%s%s]",
             safe_name, source_system, phase, agent_role, deliverable,
             f", scope={scope_item_id}" if scope_item_id else "")

    chunks = chunk(text)

    # hash the raw path (surrogatepass) so IDs stay unique even when a bad
    # byte collapsed to '?' in the display name
    doc_id    = hashlib.md5(rel_path_raw.encode("utf-8", "surrogatepass")).hexdigest()[:8]
    chunk_dir = CHUNKS_DIR / phase / agent_role
    chunk_dir.mkdir(parents=True, exist_ok=True)

    for idx, c in enumerate(chunks):
        out = chunk_dir / f"{doc_id}_{idx:04d}.json"
        out.write_text(json.dumps({
            "id":              f"{doc_id}_{idx:04d}",
            "source":          safe_name,
            "source_system":   source_system,
            "relative_path":   rel_path,
            "phase":           phase,
            "agent_role":      agent_role,
            "deliverable_type": deliverable,
            "content_type":    content_type,
            "scope_item_id":   scope_item_id,
            "client":          "[CLIENT]",
            "chunk_index":     idx,
            "total_chunks":    len(chunks),
            "text":            c,
            "ingested_at":     datetime.utcnow().isoformat() + "Z",
        }, ensure_ascii=False, indent=2), encoding="utf-8")

    log.info("  -> %d chunks saved to chunks/%s/%s/", len(chunks), phase, agent_role)
    return len(chunks)


# ── LOCAL MODE (POC — files already on EC2) ───────────────────────────────────
def process_local():
    """Process files already in brain/sharepoint/raw/ — no Graph API needed."""
    if not RAW_DIR.exists() or not any(RAW_DIR.iterdir()):
        log.error("No files found in %s — upload documents first via SCP", RAW_DIR)
        sys.exit(1)

    CHUNKS_DIR.mkdir(parents=True, exist_ok=True)
    total_chunks, total_files = 0, 0

    skipped = 0
    for f in sorted(RAW_DIR.rglob("*")):
        if not f.is_file() or f.suffix.lower() not in SUPPORTED_EXT:
            continue
        try:
            n = _ingest_one_local(f)
        except Exception as e:                 # one bad file never kills the run
            log.warning("Skipping %s: %s", _safe_str(f.name), e)
            skipped += 1
            continue
        total_chunks += n
        total_files  += 1

    if skipped:
        log.warning("Skipped %d file(s) due to errors — see warnings above.", skipped)
    log.info("Done. %d files, %d chunks across phases:", total_files, total_chunks)
    for p_dir in sorted(CHUNKS_DIR.rglob("*.json")):
        pass  # counted below
    for p_dir in sorted(CHUNKS_DIR.iterdir()):
        if p_dir.is_dir():
            count = len(list(p_dir.rglob("*.json")))
            log.info("  %-12s %d chunks", p_dir.name + ":", count)
    log.info("Next: run Bedrock Titan embeddings.")

# ── TOKEN CACHE ───────────────────────────────────────────────────────────────
def _load_cache():
    try:
        import msal
    except ImportError:
        print("Run: pip3.11 install msal requests python-docx pymupdf python-pptx")
        sys.exit(1)
    cache = msal.SerializableTokenCache()
    if TOKEN_CACHE.exists():
        cache.deserialize(TOKEN_CACHE.read_text())
    return cache

def _save_cache(cache):
    if cache.has_state_changed:
        TOKEN_CACHE.parent.mkdir(parents=True, exist_ok=True)
        TOKEN_CACHE.write_text(cache.serialize())
        TOKEN_CACHE.chmod(0o600)

# ── AUTH ──────────────────────────────────────────────────────────────────────
def get_token() -> str:
    import msal
    cache = _load_cache()
    app = msal.PublicClientApplication(
        CLIENT_ID,
        authority=f"https://login.microsoftonline.com/{TENANT_ID}",
        token_cache=cache,
    )
    result = None
    accounts = app.get_accounts()
    if accounts:
        log.info("Using cached token...")
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
    if not result:
        log.info("No cached token — starting device code flow...")
        flow = app.initiate_device_flow(scopes=SCOPES)
        if "user_code" not in flow:
            raise RuntimeError(f"Device flow error: {flow}")
        print("\n" + "=" * 60)
        print(flow["message"])
        print("=" * 60 + "\n")
        result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result.get('error_description', result)}")
    _save_cache(cache)
    log.info("Authenticated successfully.")
    return result["access_token"]

# ── GRAPH API ─────────────────────────────────────────────────────────────────
def _get(token, path, params=None):
    import requests
    resp = requests.get(
        f"{GRAPH_BASE}{path}",
        headers={"Authorization": f"Bearer {token}"},
        params=params, timeout=30,
    )
    resp.raise_for_status()
    return resp.json()

def get_site_id(token, site_url) -> str:
    hostname = site_url.split("/")[2]
    path     = "/".join(site_url.rstrip("/").split("/")[3:])
    return _get(token, f"/sites/{hostname}:/{path}")["id"]

def list_files_recursive(token, drive_id, folder_path="") -> list:
    if folder_path:
        url = f"/drives/{drive_id}/root:/{folder_path.strip('/')}:/children"
    else:
        url = f"/drives/{drive_id}/root/children"
    items = []
    while url:
        data = _get(token, url, {"$top": 200})
        for item in data.get("value", []):
            if "folder" in item:
                sub = f"{folder_path}/{item['name']}".strip("/")
                items.extend(list_files_recursive(token, drive_id, sub))
            elif "file" in item and Path(item["name"]).suffix.lower() in SUPPORTED_EXT:
                item["_folder_path"] = folder_path
                items.append(item)
        next_link = data.get("@odata.nextLink", "")
        url = next_link.replace(GRAPH_BASE, "") if next_link else None
    return items

def download(token, drive_id, item_id, dest: Path):
    import requests
    resp = requests.get(
        f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}/content",
        headers={"Authorization": f"Bearer {token}"},
        allow_redirects=True, timeout=120,
    )
    resp.raise_for_status()
    dest.write_bytes(resp.content)

# ── GRAPH API MODE ────────────────────────────────────────────────────────────
def process_graph():
    missing = [v for v in ["GRAPH_TENANT_ID","GRAPH_CLIENT_ID","SHAREPOINT_SITE_URL"]
               if not os.environ.get(v)]
    if missing:
        print(f"Missing env vars: {', '.join(missing)}")
        sys.exit(1)

    RAW_DIR.mkdir(parents=True, exist_ok=True)
    CHUNKS_DIR.mkdir(parents=True, exist_ok=True)

    token   = get_token()
    site_id = get_site_id(token, SITE_URL)

    drives = _get(token, f"/sites/{site_id}/drives")["value"]
    drive  = next((d for d in drives if d["name"] == LIBRARY), drives[0])
    did    = drive["id"]

    log.info("Listing files recursively from: %s / %s", LIBRARY, SUBFOLDER)
    files = list_files_recursive(token, did, SUBFOLDER)
    log.info("Found %d supported files", len(files))

    total_chunks, total_files = 0, 0
    for item in files:
        name        = item["name"]
        folder_path = item.get("_folder_path", "")
        bpd_scope   = detect_sap_bpd(name)
        if bpd_scope:                       # SAP standard BPD — reference, not delivery
            source_system, phase, agent_role = "sap_bpd", "Reference", "reference"
            deliverable, scope_item_id = "business_process_doc", bpd_scope
        else:                               # client delivery document
            source_system, scope_item_id = "sharepoint", None
            phase       = detect_phase(folder_path)
            agent_role  = detect_agent_role(folder_path)
            deliverable = detect_deliverable_type(folder_path)
        content_type = detect_content_type(name)
        dest        = RAW_DIR / name

        log.info("Downloading: %s [src=%s, phase=%s, agent=%s%s]",
                 name, source_system, phase, agent_role,
                 f", scope={scope_item_id}" if scope_item_id else "")
        download(token, did, item["id"], dest)

        text   = extract_text(dest)
        text   = mask(text)
        chunks = chunk(text)

        doc_id    = hashlib.md5(f"{folder_path}/{name}".encode()).hexdigest()[:8]
        chunk_dir = CHUNKS_DIR / phase / agent_role
        chunk_dir.mkdir(parents=True, exist_ok=True)

        for idx, c in enumerate(chunks):
            out = chunk_dir / f"{doc_id}_{idx:04d}.json"
            out.write_text(json.dumps({
                "id":              f"{doc_id}_{idx:04d}",
                "source":          name,
                "source_system":   source_system,
                "folder_path":     folder_path,
                "phase":           phase,
                "agent_role":      agent_role,
                "deliverable_type": deliverable,
                "content_type":    content_type,
                "scope_item_id":   scope_item_id,
                "client":          "[CLIENT]",
                "chunk_index":     idx,
                "total_chunks":    len(chunks),
                "text":            c,
                "ingested_at":     datetime.utcnow().isoformat() + "Z",
            }, ensure_ascii=False, indent=2))

        total_chunks += len(chunks)
        total_files  += 1
        log.info("  -> %d chunks [phase=%s, agent=%s]", len(chunks), phase, agent_role)

    log.info("Done. %d files, %d total chunks.", total_files, total_chunks)
    for p_dir in sorted(CHUNKS_DIR.iterdir()):
        if p_dir.is_dir():
            count = len(list(p_dir.rglob("*.json")))
            log.info("  %-12s %d chunks", p_dir.name + ":", count)

# ── CLEAN ─────────────────────────────────────────────────────────────────────
def clean_chunks():
    """Wipe the chunks output dir so a re-run has no stale/mis-classified chunks.
    Never touches raw/ (the uploaded source documents are preserved)."""
    if CHUNKS_DIR.exists():
        shutil.rmtree(CHUNKS_DIR)
        log.info("Cleaned chunks dir: %s", CHUNKS_DIR)
    CHUNKS_DIR.mkdir(parents=True, exist_ok=True)

# ── ENTRY POINT ───────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--local", action="store_true",
                        help="Process files already in brain/sharepoint/raw/ (POC mode)")
    parser.add_argument("--clean", action="store_true",
                        help="Wipe chunks/ before ingest (avoids stale chunks from "
                             "earlier runs with old masking/classification). raw/ is kept.")
    args = parser.parse_args()

    if args.clean:
        clean_chunks()

    if args.local:
        log.info("Running in LOCAL mode — processing files from %s", RAW_DIR)
        process_local()
    else:
        log.info("Running in GRAPH API mode")
        process_graph()

if __name__ == "__main__":
    main()
